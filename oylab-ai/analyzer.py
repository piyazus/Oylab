import io, os, re, json, requests
from typing import Dict, List, BinaryIO

import pdfplumber
from pptx import Presentation

# ---------- Weights ----------
WEIGHTS = {
    "team": 0.25,
    "market": 0.20,
    "product": 0.30,
    "finance": 0.20,
    "design": 0.05,
}

# ---------- Keyword sets ----------
TEAM = [r"\bteam\b", r"\bfounder\b", r"\bco[- ]?founder\b", r"\bcto\b", r"\bceo\b",
        r"\bexperience\b", r"\baward\b", r"\badvisor\b", r"\bmentor\b"]
MARKET = [r"\bmarket\b", r"\btam\b", r"\bsam\b", r"\bsom\b", r"\busers?\b",
          r"\bgrowth\b", r"\bcustomer\b", r"\bsegment\b", r"\btraction\b",
          r"\bcampaign\b", r"\binfluencer\b", r"\bpilot\b"]
PRODUCT = [r"\bproduct\b", r"\bproblem\b", r"\bsolution\b", r"\bmvp\b",
           r"\bprototype\b", r"\btech(nology)?\b", r"\barchitecture\b",
           r"\balgorithm\b", r"\broadmap\b", r"\bapi\b", r"\bmobile app\b",
           r"\bbackend\b", r"\bfrontend\b"]
FINANCE = [r"\brevenue\b", r"\bpricing\b", r"\bcost\b", r"\bunit\b",
           r"\bcogs\b", r"\bcac\b", r"\bltv\b", r"\bmargin\b",
           r"\bmoneti[sz]ation\b", r"\bgtm\b"]
DESIGN = [r"\bdesign\b", r"\bui\b", r"\bux\b", r"\bmockups?\b", r"\bwireframes?\b",
          r"\bfigma\b", r"\bprototype\b", r"\bvisuals?\b", r"\btypography\b",
          r"\blayout\b", r"\bstyle\b", r"\bbrand(ing)?\b"]

# ---------- Extraction ----------
def _extract_pdf(buf: BinaryIO) -> str:
    text = []
    with pdfplumber.open(buf) as pdf:
        for p in pdf.pages:
            t = p.extract_text() or ""
            text.append(t)
    return "\n".join(text)

def _extract_ppt(buf: BinaryIO) -> str:
    prs = Presentation(buf)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                chunks.append(shape.text)
    return "\n".join(chunks)

def _extract_any(file_storage) -> str:
    name = (getattr(file_storage, "filename", "") or "").lower()
    raw = file_storage.read()
    file_storage.stream.seek(0)
    buf = io.BytesIO(raw)
    if name.endswith(".pdf"):
        return _extract_pdf(buf)
    if name.endswith(".pptx") or name.endswith(".ppt"):
        return _extract_ppt(buf)
    try:
        buf.seek(0);  return _extract_pdf(buf)
    except Exception:
        buf.seek(0);  return _extract_ppt(buf)

# ---------- Heuristic ----------
def _score(text: str, patterns: List[str]) -> int:
    t = text.lower()
    found = 0
    for p in patterns:
        m = re.findall(p, t)
        if m:
            found += 1
            if len(m) > 2:  # бонус за частоту
                found += 1
    maxp = max(1, len(patterns) + 2)
    return int((found / maxp) * 100)

def _recommend(b: Dict[str, int]) -> List[str]:
    recs: List[str] = []
    if b["team"] < 70:
        recs.append("Add a team slide: roles, achievements, why this team wins.")
    if b["market"] < 70:
        recs.append("Quantify TAM/SAM/SOM and add customer validation/segments.")
    if b["product"] < 70:
        recs.append("Clarify problem→solution; show MVP screenshots or demo link.")
    if b["finance"] < 70:
        recs.append("Explain pricing, unit economics (CAC/LTV, margin) and GTM.")
    if b["design"] < 70:
        recs.append("Improve visuals: consistent UI/UX, Figma mockups, clear layout.")
    if not recs:
        recs.append("Great fundamentals. Add traction metrics and a clear ‘ask’.")
    return recs

def _heuristic_breakdown(text: str) -> Dict[str, int]:
    return {
        "team": _score(text, TEAM),
        "market": _score(text, MARKET),
        "product": _score(text, PRODUCT),
        "finance": _score(text, FINANCE),
        "design": _score(text, DESIGN),
    }

# ---------- OpenRouter LLM ----------
def _clamp(n: int) -> int:
    return max(0, min(100, int(n)))

def _llm_analyze_openrouter(text: str) -> dict:
    url = "https://openrouter.ai/api/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {os.environ['OPENAI_API_KEY']}",
        "Content-Type": "application/json"
    }
    if len(text) > 60000:
        text = text[:60000]

    payload = {
        "model": os.getenv("OPENAI_MODEL", "openai/gpt-4o-mini"),
        "response_format": {"type": "json_object"},
        "messages": [
            {"role": "system", "content": (
                "You are a strict pitch-deck evaluator for high-school startup competitions. "
                "Score each category from 0 to 100 based on evidence in the text. "
                "Never guess; if information is missing, give a low score. Output compact JSON only."
            )},
            {"role": "user", "content": (
                "Return a JSON object with keys: scores, suggestions.\n"
                "scores = {team:int, market:int, product:int, finance:int, design:int}.\n"
                "suggestions = array of 3-8 short actionable strings.\n\n"
                f"<deck>\n{text}\n</deck>"
            )}
        ]
    }

    r = requests.post(url, headers=headers, json=payload, timeout=60)
    r.raise_for_status()
    data = r.json()
    content = data["choices"][0]["message"]["content"]
    parsed = json.loads(content)

    scores = parsed.get("scores", {})
    breakdown = {
        "team": _clamp(scores.get("team", 0)),
        "market": _clamp(scores.get("market", 0)),
        "product": _clamp(scores.get("product", 0)),
        "finance": _clamp(scores.get("finance", 0)),
        "design": _clamp(scores.get("design", 0)),
    }
    score = int(sum(breakdown[k] * WEIGHTS[k] for k in WEIGHTS))
    return {"score": score, "breakdown": breakdown, "recommendations": parsed.get("suggestions", [])}

# ---------- Orchestrator ----------
def analyze_pitch(file_storage) -> Dict:
    text = _extract_any(file_storage)
    mode = os.getenv("AI_MODE", "heuristic").lower()
    try:
        if mode == "openai" and os.getenv("OPENAI_API_KEY"):
            return _llm_analyze_openrouter(text)
    except Exception as e:
        print("LLM failed, fallback to heuristic:", e)
    # fallback
    b = _heuristic_breakdown(text)
    score = int(sum(b[k] * WEIGHTS[k] for k in WEIGHTS))
    return {"score": score, "breakdown": b, "recommendations": _recommend(b)}
